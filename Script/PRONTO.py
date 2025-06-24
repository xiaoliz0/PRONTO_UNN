#!/usr/bin/env python
# -*- coding: utf-8 -*-
#
# This script needs to be run with python above version 3.
# To install module xlutils, run the command: sudo pip install xlutils
# To install module python-pptx, run the command: sudp pip3 install python-pptx
"""
This script is a tool used to filter and analysis data from TSO500 results.
And generate the PP report based on the results data and template file.
"""

import os
import re
import sys
import getopt
import time
from configparser import ConfigParser
from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE

DNA_sampleID = ""
RNA_sampleID = ""
runID = ""
batch_nr = ""
tumor_content_nr = ""
ipd_birth_year = ""
ipd_diagnosis_year = "-"
ipd_age = ""
ipd_gender = ""
ipd_consent = ""
ipd_collection_year = "-"
requisition_hospital = ""
pathology_comment = ""
ipd_material_id = ""
DNA_material_id = ""
RNA_material_id = ""
sample_info_comment = ""
extraction_hospital = ""
inclusion_site = ""
ipd_clinical_diagnosis = "-"
sample_material = ""
sample_type = ""
tumor_type = ""
MSI_TSO500 = ""

# Read clinical data from MTF files version from 2025.
def get_patient_info_from_MTF(ipd_material_file,ipd_no,DNA_sampleID):
	import xlrd
	global RNA_sampleID
	global ipd_birth_year
	global ipd_age
	global ipd_clinical_diagnosis
	global ipd_gender
	global ipd_consent
	global DNA_material_id
	global RNA_material_id
	global ipd_collection_year
	global requisition_hospital
	global pathology_comment
	global sample_info_comment
	global extraction_hospital
	global batch_nr
	global tumor_content_nr
	global inclusion_site
	open_exl_material = xlrd.open_workbook(ipd_material_file)
	sheet_material = open_exl_material.sheet_by_index(0)
	nrows_material = sheet_material.nrows
	ncols_material = sheet_material.ncols
	columns = {
		'ipd': 'InPreD ID',
		'gender': 'Gender',
		'age': 'Age',
		'birth_date': 'Date of birth',
		'DIT_number': 'DIT number',
		'consent': 'Study ID',
		'requisition_hospital': 'Requester Hospital',
		'Histopathological_diagnosis': 'diagnosis',
		'comment': 'Comments',
		'material_id': 'Sample material ID',
		'tumor_content_nr': 'Tumor cells [%]',
		'sample_ID': 'Sample ID',
		'ex_pathology_info': 'Molecular Pathology information',
		'ex_sample_info': 'Sample information',
		'ex_data_section': 'Extraction Data',
		'ex_library_pre': 'Library Preparation (LP) Data',
		'extraction_hospital': 'Extraction Hospital',
		'batch_nr': 'LP batch'
		}
	ipd_birth_date = ""
	sample_info_row = 0
	extra_data_row = 0
	library_pre_row = 0
	for l in range(nrows_material):
		if(sheet_material.cell_value(l,0) == columns['ex_sample_info']):
			sample_info_row = l
		if(sheet_material.cell_value(l,0) == columns['ex_data_section']):
			extra_data_row = l
		if(sheet_material.cell_value(l,0) == columns['ex_library_pre']):
			library_pre_row = l  
	for r in range(nrows_material):
		for c in range(ncols_material):
			if(sheet_material.cell_value(r,c) == columns['ipd']):
				ipd_MTF = sheet_material.cell_value(r+2,c)
				if(ipd_MTF != ipd_no):
					print("""               Error:
                        The InPreD patient ID in IPD Material Transit Form InPreD NGS file does not match with the IPD number! 
                        Please check and fix the mistake before run this script again!""")
					print("                 IPD is " + ipd_MTF + " in MTF, while IPD is " + ipd_no[3:] + " in TSO500.")
					sys.exit(0)
			if(sheet_material.cell_value(r,c) == columns['birth_date']):
				ipd_birth_date_exl = sheet_material.cell_value(r+2,c)
				try:
					datetime_date = str(xlrd.xldate_as_datetime(ipd_birth_date_exl,0))
					ipd_birth_year = datetime_date.split('-')[0]
				except:
					ipd_birth_year = "-"
			if(sheet_material.cell_value(r,c) == columns['gender'] and ipd_gender == ""):
				ipd_gender = str(sheet_material.cell_value(r+2,c))
			if(sheet_material.cell_value(r,c) == columns['age']):
				ipd_age = str(sheet_material.cell_value(r+2,c))
			if(sheet_material.cell_value(r,c) == columns['Histopathological_diagnosis'] and sheet_material.cell_value(r+1,c) != ""):
				ipd_clinical_diagnosis = str(sheet_material.cell_value(r+1,c))
			if(sheet_material.cell_value(r,c) == columns['consent'] and ipd_consent == ""):
				ipd_consent = str(sheet_material.cell_value(r+2,c))
				if(ipd_consent == "0.0"):
					ipd_consent = ""
				for r in range(r,(sample_info_row-2)):
					DIT_number = "-"
					comments = "-"
					if(sheet_material.cell_value(r,0) == columns['DIT_number'] and sheet_material.cell_value(r+2,0) != ""):
						DIT_number = sheet_material.cell_value(r+2,0)
					if(sheet_material.cell_value(r,6) == columns['requisition_hospital'] and requisition_hospital == ""):
						requisition_hospital = sheet_material.cell_value(r+2,6)
					if((sheet_material.cell_value(r+2,c) != "" and sheet_material.cell_value(r+2,c) != "-" and str(sheet_material.cell_value(r+2,c)) != "0.0") and str(sheet_material.cell_value(r+2,c)) not in ipd_consent):
						if(ipd_consent == ""):
							ipd_consent = str(sheet_material.cell_value(r+2,c))
						else:
							ipd_consent = ipd_consent + "," + str(sheet_material.cell_value(r+2,c))
					if(sheet_material.cell_value(r,10) == columns['comment'] and sheet_material.cell_value(r+2,10) != "" and str(sheet_material.cell_value(r+2,10)) != "0.0" and str(sheet_material.cell_value(r+2,10)) != "0"):
						comments = str(sheet_material.cell_value(r+2,10))
					if(pathology_comment == ""):
						pathology_comment = DIT_number + ":" + comments
					else:
						if(DIT_number != "-" or comments != "-"):
							pathology_comment += "|" + DIT_number + ":" + comments
			if(sheet_material.cell_value(r,c) == columns['material_id'] and ipd_material_id == ""):
				for r in range(r,(extra_data_row-2)):
					sample_ID = "-"
					comments = "-"
					if(sheet_material.cell_value(r+2,9) == DNA_sampleID and sheet_material.cell_value(r+2,c) != "" and str(sheet_material.cell_value(r+2,c)) not in DNA_material_id):
						if(DNA_material_id == ""):
							DNA_material_id = str(sheet_material.cell_value(r+2,c))
						else:
							DNA_material_id = DNA_material_id + "," + str(sheet_material.cell_value(r+2,c))
						tumor_content_nr = sheet_material.cell_value(r+2,2)
					if(sheet_material.cell_value(r+2,c) != "" and sheet_material.cell_value(r+2,9)[8] == 'R' and str(sheet_material.cell_value(r+2,c)) not in RNA_material_id):
						RNA_sampleID = str(sheet_material.cell_value(r+2,9)) 
						if(RNA_material_id == ""):
							RNA_material_id = str(sheet_material.cell_value(r+2,c))
						else:
							RNA_material_id = RNA_material_id + "," + str(sheet_material.cell_value(r+2,c))
					if(sheet_material.cell_value(r+2,9) != ""):
						sample_ID = sheet_material.cell_value(r+2,9)
					if(sheet_material.cell_value(r+2,10) != "" and str(sheet_material.cell_value(r+2,10)) != "0.0" and str(sheet_material.cell_value(r+2,10)) != "0"):
						comments = str(sheet_material.cell_value(r+2,10))
					if(sample_info_comment == ""):
						sample_info_comment = "{}: {}".format(sample_ID, comments)
					else:
						if(sample_ID != "-" or comments != "-"):
							sample_info_comment += "|" + sample_ID + ": " + comments
			if(sheet_material.cell_value(r,c) == columns['extraction_hospital'] and extraction_hospital == ""):
				for r in range(r,(library_pre_row-2)):
					if(sheet_material.cell_value(r+2,8) == DNA_sampleID):
						extraction_hospital = str(sheet_material.cell_value(r+2,c))
						break
			if(sheet_material.cell_value(r,c) == columns['batch_nr'] and batch_nr == ""):
				for r in range(r,(nrows_material-2)):
					if(sheet_material.cell_value(r+2,0) == DNA_sampleID):
						batch_nr = str(sheet_material.cell_value(r+2,c))
	open_exl_material.release_resources()
	if(ipd_consent == "0.0"):
		ipd_consent = ""
	if(ipd_age == "" and ipd_birth_date != ""):
		ipd_age = "<1"
	inclusion_site_list = {'R': 'Radium', 'U': 'Ullevål', 'C': 'Riksen', 'A': 'Ahus', 'D': 'Drammen', 'B': 'Bærum', 'G': 'Gjøvik', 'I': 'Hamar', 'L': 'Lillehammer', 'T': 'Vestfold', 'K': 'Sørlandet', 'Q': 'Østfold', 'V': 'Telemark', 'Y': 'Lovisenberg', 'H': 'Haukeland', 'S': 'Stavanger', 'E': 'Fonna', 'F': 'Førde', 'O': 'St.Olavs', 'M': 'Nord-trøndelag', 'J': 'Møre og Romsdal', 'N': 'Nord Norge', 'P': 'Nordland'}
	if("IKKE IMPRESS" in ipd_consent):
		inclusion_site = ""
	else:
		try:
			site_letter_code = ipd_consent[-6]
			inclusion_site = inclusion_site_list.get(site_letter_code)
		except:
			inclusion_site = "Inclusion site"


def get_RNA_material_id(InPreD_clinical_data_file,RNA_sampleID,encoding_sys):
	RNA_material_id_exist = False
	if(encoding_sys != ""):
		f = open(InPreD_clinical_data_file, 'r', encoding=encoding_sys)
	else:
		f = open(InPreD_clinical_data_file, 'r')
	for l in f:
		if(RNA_sampleID == l.split('\t')[0]):
			RNA_material_id = l.split('\t')[8]
			RNA_material_id_exist = True
	f.close()
	if(RNA_material_id_exist == False):
		print("Warning: The "+ RNA_sampleID + " does not exist in the meta file! The report will be generated without RNA sample material ID!")
		RNA_material_id = ""
	return RNA_material_id


def update_ppt_template_data(inpred_node,ipd_no,ipd_gender,ipd_age,ipd_diagnosis_year,DNA_material_id,RNA_material_id,ipd_consent,requisition_hospital,pathology_comment,ipd_clinical_diagnosis,tumor_type,sample_type,sample_material,sample_info_comment,tumor_content,ppt_template,output_ppt_file):
	if(ipd_gender != "" and ipd_gender != "X"):
		gender = ipd_gender[0]
	else:
		gender = ""
	if(ipd_age != "" and ipd_age != "-" and ipd_age != "XX" and ipd_age != "<1"):
		age = str(int(float(ipd_age)))
	else:
		age = ipd_age
	try:
		sample = sample_type + '\n' + sample_material
	except:
		sample = ""
	today_date = time.strftime("%d", time.localtime())
	today_month = time.strftime("%b", time.localtime())
	today_year = time.strftime("%Y", time.localtime())
	today = today_date + '\n' + today_month.upper() + '\n' + today_year
	ppt = Presentation(ppt_template)
	indexs = [1,3,4,5,6]
	for index in indexs:
		slide = ppt.slides[index]
		textbox1 = slide.shapes.add_textbox(Inches(3.75), Inches(0.11), Inches(1.33), Inches(0.50))
		tf1 = textbox1.text_frame
		tf1.paragraphs[0].text = ipd_no
		tf1.paragraphs[0].font.size = Pt(24)
		tf1.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		tf1.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox2 = slide.shapes.add_textbox(Inches(8.99), Inches(0.02), Inches(0.45), Inches(0.55))
		tf2 = textbox2.text_frame
		tf2.paragraphs[0].text = today
		tf2.paragraphs[0].font.size = Pt(9)
		tf2.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
		tf2.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM
		textbox4 = slide.shapes.add_textbox(Inches(0.50), Inches(1.47), Inches(0.86), Inches(0.25))
		tf4 = textbox4.text_frame
		tf4.paragraphs[0].text = requisition_hospital
		tf4.paragraphs[0].font.size = Pt(10)
		tf4.paragraphs[0].alignment = PP_ALIGN.CENTER
		tf4.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		textbox5 = slide.shapes.add_textbox(Inches(0.71), Inches(1.84), Inches(0.86), Inches(0.50))
		tf5 = textbox5.text_frame
		tf5.paragraphs[0].text = sample
		tf5.paragraphs[0].font.size = Pt(8)
		tf5.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox6 = slide.shapes.add_textbox(Inches(0.81), Inches(2.65), Inches(0.63), Inches(0.33))
		tf6 = textbox6.text_frame
		tf6.paragraphs[0].text = tumor_content
		tf6.paragraphs[0].font.size = Pt(14)
		tf6.paragraphs[0].alignment = PP_ALIGN.CENTER
		if(index == 1 or ipd_clinical_diagnosis == "-" or ipd_clinical_diagnosis == ""):
			textbox7 = slide.shapes.add_textbox(Inches(5.77), Inches(0.19), Inches(0.86), Inches(0.33))
			tf7 = textbox7.text_frame
			tf7.paragraphs[0].text = str(tumor_type)
			tf7.paragraphs[0].font.size = Pt(14)
		else:
			if('\n' in ipd_clinical_diagnosis):
				textbox7 = slide.shapes.add_textbox(Inches(4.95), Inches(0.11), Inches(2.36), Inches(0.50))
				tf7 = textbox7.text_frame
				tf7.paragraphs[0].font.size = Pt(12)
			else:
				textbox7 = slide.shapes.add_textbox(Inches(5.77), Inches(0.19), Inches(0.86), Inches(0.33))
				tf7 = textbox7.text_frame
				tf7.paragraphs[0].font.size = Pt(14)
			tf7.paragraphs[0].text = ipd_clinical_diagnosis
		tf7.paragraphs[0].font.italic = True
		tf7.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		tf7.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox11 = slide.shapes.add_textbox(Inches(0.85), Inches(1.12), Inches(0.48), Inches(0.27))
		tf11 = textbox11.text_frame
		tf11.paragraphs[0].text = ipd_diagnosis_year
		tf11.paragraphs[0].font.size = Pt(10)
		tf11.paragraphs[0].alignment = PP_ALIGN.LEFT
		textbox12 = slide.shapes.add_textbox(Inches(0.61), Inches(0.35), Inches(1.02), Inches(0.33))
		tf12 = textbox12.text_frame
		tf12.paragraphs[0].text = inpred_node
		tf12.paragraphs[0].font.size = Pt(14)
		tf12.paragraphs[0].alignment = PP_ALIGN.CENTER
		tf12.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		if(index == 1):
			gender_age = ""
			ipd_material_id_index = ""
			ipd_consent_index = ""
		textbox8 = slide.shapes.add_textbox(Inches(0.69), Inches(0.79), Inches(0.87), Inches(0.40))
		tf8 = textbox8.text_frame
		tf8.paragraphs[0].text = gender_age
		tf8.paragraphs[0].font.size = Pt(18)
		tf8.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox9 = slide.shapes.add_textbox(Inches(0.73), Inches(2.25), Inches(0.70), Inches(0.26))
		tf9 = textbox9.text_frame
		tf9.paragraphs[0].text = ipd_material_id_index
		tf9.paragraphs[0].font.size = Pt(5)
		tf9.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox10 = slide.shapes.add_textbox(Inches(2.10), Inches(0.11), Inches(1.07), Inches(0.50))
		tf10 = textbox10.text_frame
		tf10.paragraphs[0].text = ipd_consent_index
		tf10.paragraphs[0].font.size = Pt(14)
		tf10.paragraphs[0].alignment = PP_ALIGN.CENTER
		tf10.paragraphs[0].font.italic = True
		tf10.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		if(index == 3):
			textbox11 = slide.shapes.add_textbox(Inches(1.85), Inches(1.25), Inches(3.25), Inches(0.27))
			tf11 = textbox11.text_frame
			tf11.paragraphs[0].text = pathology_comment + "\n\n" + sample_info_comment.replace("|","\n")
			tf11.paragraphs[0].font.size = Pt(10)
			tf11.paragraphs[0].alignment = PP_ALIGN.LEFT
		gender_age = gender + '/' + age + 'y'
		if(RNA_material_id != ""):
			ipd_material_id_index = "DNA:" + DNA_material_id + "\nRNA:" + RNA_material_id
		else:
			ipd_material_id_index = "DNA:" + DNA_material_id
		ipd_consent_index = "Trial ID\n" + ipd_consent 

	ppt.save(output_ppt_file)


def update_clinical_master_file(InPreD_clinical_data_file,sample_id,if_generate_report,ipd_birth_year,clinical_diagnosis,ipd_gender,ipd_consent,material_id,ipd_collection_year,requisition_hospital,extraction_hospital,tumor_content_nr,batch_nr,pathology_comment,sample_info_comment,encoding_sys):
	global ipd_diagnosis_year
	global runID
	if_exist = False
	new_content = ""
	if(encoding_sys != ""):
		fr = open(InPreD_clinical_data_file, 'r', encoding=encoding_sys)
	else:
		fr = open(InPreD_clinical_data_file, 'r')
	if(str(requisition_hospital) == "0.0"):
		requisition_hospital = "-"
	if(str(extraction_hospital) == "0.0"):
		extraction_hospital = "-"
	for ln in fr:
		if(ln.split('\t')[0] == sample_id):
			if_exist = True
			line = '\t'.join([sample_id, runID, if_generate_report, ipd_birth_year, ipd_diagnosis_year, clinical_diagnosis, ipd_gender[0], ipd_consent, material_id, ipd_collection_year, requisition_hospital, extraction_hospital, str(tumor_content_nr), batch_nr, pathology_comment, sample_info_comment + '\n'])
			new_line = ln.replace(ln,line)
			new_content = new_content + new_line
		else:
			new_content = new_content + ln
	fr.close()
	if(encoding_sys != ""):
		fa = open(InPreD_clinical_data_file, 'a', encoding=encoding_sys)
	else:
		fa = open(InPreD_clinical_data_file, 'a')
	if(if_exist == False):
		line = '\t'.join([sample_id, runID, if_generate_report, ipd_birth_year, ipd_diagnosis_year, clinical_diagnosis, ipd_gender[0], ipd_consent, material_id, ipd_collection_year, requisition_hospital, extraction_hospital, str(tumor_content_nr), batch_nr, pathology_comment, sample_info_comment + '\n'])
		if(encoding_sys != ""):
			fa = open(InPreD_clinical_data_file, 'a', encoding=encoding_sys)
		else:
			fa = open(InPreD_clinical_data_file, 'a')
		fa.write(line)
		fa.close()
	else:
		if(encoding_sys != ""):
			fw = open(InPreD_clinical_data_file, 'w', encoding=encoding_sys)
		else:
			fw = open(InPreD_clinical_data_file, 'w')
		fw.write(new_content)
		fw.close()


def usage(exit_status = 0):
	print ("""Usage: python3  %s
        This script is a tool used to generate the paitent report based on the TSO500 analysis results and the personal intomation from the clinical data in In/InPreD_PRONTO_metadata.txt,
	and update the SOPPI results into the file Out/InPreD_PRONTO_metadata_tsoppi.txt when the reports are generated.
	This script could also fill the patient personal information into the clinical data file with the MTF files under the foder In/MTF/. (This fuction currently is only used by OUS)
	To run this script tool in your computer with python3, it will read the clinical data from In/InPreD_PRONTO_metadata.txt and generate reports for the Sample_id with Create_report==Y:
	
	python3 InPreD_PRONTO.py

	Extra parameters:
	-c, --clinical_file Fill the patient personal information into the clinical data file: InPreD_PRONTO_metadata.txt with the MTF files under the foder In/MTF/
	python3 InPreD_PRONTO.py -D <DNA_sampleID> -c
	or:
	python3 InPreD_PRONTO.py --DNAsampleID=<DNA_sampleID> --clinicalFile

	This script will create sub-folder with runID/IPDXXX in Out/, move the IPD_Material file into it and generate all the results files under that sub-folder.
	 
	-h, --help See this help information and exit.
        """ % sys.argv[0])
	sys.exit(exit_status)


def main(argv):
	global DNA_sampleID
	global RNA_sampleID
	global ipd_birth_year
	global ipd_clinical_diagnosis
	global ipd_gender
	global ipd_consent
	global ipd_material_id
	global DNA_material_id
	global RNA_material_id
	global ipd_collection_year
	global pathology_comment
	global requisition_hospital
	global extraction_hospital
	global sample_material
	global sample_type
	global tumor_type
	global batch_nr
	global tumor_content_nr
	global sample_info_comment
	update_clinical_file = False
	tumor_content = "XX"
	ipd_material_file_new = ""
	try:
		opts, args = getopt.getopt(sys.argv[1:], "hD:c", ["help", "DNAsampleID=", "clinicalFile"])
	except getopt.GetoptError:
		usage(1)

	for opt, arg in opts:
		if opt in ("-h", "--help"):
			usage()
		elif opt in ("-D", "--DNAsampleID"):
			DNA_sampleID = arg
			ipd_no = DNA_sampleID.split('-')[0]
		elif opt in ("-c", "--clinicalFile"):
			update_clinical_file = True
	DNA_sampleID_format = '^IP[A-Z]\\d{4}-D(\\d|X){2}-[A-z](\\d|X){2}-[A-z](\\d|X){2}'

	base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
	config_file = base_dir + "/Config/configure_PRONTO.ini"
	InPreD_clinical_data_file = base_dir + "/In/InPreD_PRONTO_metadata.txt"
	output_path_root = base_dir + "/Out/"
	InPreD_clinical_tsoppi_data_file = base_dir + "/Out/InPreD_PRONTO_metadata_tsoppi.txt"
	cfg = ConfigParser()
	cfg.read(config_file)
	inpred_node = cfg.get("INPUT", "inpred_node")
	encoding_sys = cfg.get("INPUT", "encoding_sys")

	if not os.path.exists(InPreD_clinical_data_file):
		print ("""      Error:
		The InPreD clinical file InPreD_PRONTO_metadata.txt does not exist!
		""")
		sys.exit(0)
	if(update_clinical_file == True):
		ipd_material_file = base_dir + "/In/MTF/" + ipd_no[:3] + '-' + ipd_no[3:] + "_Material Transit Form InPreD NGS_2025.xlsx"
		if not(os.path.exists(ipd_material_file)):
			print ("""Error: IPD Material Transit Form InPreD NGS file does not exit under the MTF dir. PRONTO meta file could not be updated with patient personal information by parameter -c of this script!""")
			sys.exit(0)
		if not(re.fullmatch(DNA_sampleID_format, DNA_sampleID)):
			print("Warning: " + DNA_sampleID + " does not fit for the sample id format!")
		else:
			if os.path.exists(ipd_material_file):
				get_patient_info_from_MTF(ipd_material_file,ipd_no,DNA_sampleID)
			if_generate_report = "Y"
			update_clinical_master_file(InPreD_clinical_data_file,DNA_sampleID,if_generate_report,ipd_birth_year,ipd_clinical_diagnosis,ipd_gender,ipd_consent,DNA_material_id,ipd_collection_year,requisition_hospital,extraction_hospital,tumor_content_nr,batch_nr,pathology_comment,sample_info_comment,encoding_sys)
			print("Clinical data is added into PRONTO meta file for sample: " + DNA_sampleID)
			if(RNA_sampleID != ""):
				if_generate_report = "-"
				update_clinical_master_file(InPreD_clinical_data_file,RNA_sampleID,if_generate_report,ipd_birth_year,ipd_clinical_diagnosis,ipd_gender,ipd_consent,RNA_material_id,ipd_collection_year,requisition_hospital,extraction_hospital,tumor_content_nr,batch_nr,pathology_comment,sample_info_comment,encoding_sys)
				print("Clinical data is added into PRONTO meta file for sample: " + RNA_sampleID)
			sys.exit(0)
	ppt_nr = 0
	if(encoding_sys != ""):
		meta_file = open(InPreD_clinical_data_file, encoding=encoding_sys)
	else:
		meta_file = open(InPreD_clinical_data_file)
	with meta_file as f:
		lines = f.readlines()
		for i in range(len(lines)-1):
			ln = lines[i]
			if not(ln.startswith("#") or ln == ""):
				if(ln.split('\t')[2] == "Y"):
					ln = ln.replace("\n", "")
					ln = ln + '\t' * (15-ln.count('\t'))
					[DNA_sampleID, runID_DNA, _, ipd_birth_year, ipd_diagnosis_year, ipd_clinical_diagnosis_meta, ipd_gender, ipd_consent, DNA_material_id, ipd_collection_year, requisition_hospital, extraction_hospital, tumor_content_nr, batch_nr, pathology_comment, sample_info_comment] = ln.split('\t')
					if not(re.fullmatch(DNA_sampleID_format, DNA_sampleID)):
						print("Warning: " + DNA_sampleID + " does not fit for the sample id format!")
					try:
						ipd_age = str(int(time.strftime("%Y", time.localtime())) - int(ipd_birth_year))
					except:
						ipd_age = "-"
					ipd_no = DNA_sampleID.split('-')[0]
					next_line = lines[i+1].strip()
					if(next_line.split('\t')[0].startswith(ipd_no+"-R")):
						RNA_sampleID = next_line.split('\t')[0]

					try:
						ipd_clinical_diagnosis_ppt = ipd_clinical_diagnosis_meta.split("(")[0] + "\n(" + ipd_clinical_diagnosis_meta.split("(")[1]
					except:
						ipd_clinical_diagnosis_ppt = ipd_clinical_diagnosis_meta
					try:
						tumor_content = '~' + str(int(float(tumor_content_nr))) + '%'
					except:
						tumor_content = "XX"
					output_path = output_path_root + "/" + DNA_sampleID + "/"
					output_file_preMTB_table_path = output_path + DNA_sampleID
					if not os.path.exists(output_path):
						os.makedirs(output_path)
					ppt_template = base_dir + "/In/Template/InPreD_MTB_template.pptx"
					output_ppt_file = output_path + DNA_sampleID + "_MTB_report.pptx"
					try:
						sample_type_string = DNA_sampleID.split('-')[2]
						sample_type_short = sample_type_string[0:1]
						sample_type_list = {'M': 'Metastasis', 'T': 'Primary Tumor', 'C': 'Cell-line', 'N': 'Normal/Control', 'P': 'Primary tumor\n naive', 'p': 'Primary tumor\n post-treatment', 'R': 'Regional met\n naive', 'r': 'Regional met\n post-treatment', 'D': 'Distal met\n naive', 'd': 'Distal met\n post-treatment', 'L': 'Liquid', 'E': 'naive', 'e': 'post treatment', 'A': 'post allo transplant', 'X': 'Unknown'}
						sample_type = sample_type_list.get(sample_type_short)
					except:
						sample_type = ""
					try:
						sample_material_string = DNA_sampleID.split('-')[3]
						sample_material_short = sample_material_string[0:1]
						sample_material_list = {'F': 'Fresh Frozen', 'A': 'Archived FFPE', 'B': 'Blood', 'C': 'Cytology', 'M': 'Fresh bone marrow',  'E': 'Extramedullary','S': 'Buccal swab (normal)', 'X': 'Unspecified'}
						sample_material = sample_material_list.get(sample_material_short)
					except:
						sample_material = ""
					try:
						tumor_type_no = sample_material_string[1:3]
						tumor_type_list = {'00': 'Cancer origo incerta', '01': 'Adrenal Gland', '02': 'Ampulla of Vater', '03': 'Biliary Tract', '04': 'Bladder/Urinary Tract', '05': 'Bone', '06': 'Breast', '07': 'Cervix', '08': 'CNS/Brain', '09': 'Colon/Rectum', '10': 'Esophagus/Stomach', '11': 'Eye', '12': 'Head and Neck', '13': 'Kidney', '14': 'Liver', '15': 'Lung', '16': 'Lymphoid', '17': 'Myeloid', '18': 'Ovary/Fallopian Tube', '19': 'Pancreas', '20': 'Peripheral Nervous System', '21': 'Peritoneum', '22': 'Pleura', '23': 'Prostate', '24': 'Skin', '25': 'Soft Tissue', '26': 'Testis', '27': 'Thymus', '28': 'Thyroid', '29': 'Uterus', '30': 'Vulva/Vagina', 'XX': 'Not available'}
						tumor_type = tumor_type_list.get(tumor_type_no)
					except:
						tumor_type = ""
                
					RNA_material_id = get_RNA_material_id(InPreD_clinical_data_file,RNA_sampleID,encoding_sys)
					update_ppt_template_data(inpred_node,ipd_no,ipd_gender,ipd_age,ipd_diagnosis_year,DNA_material_id,RNA_material_id,ipd_consent,requisition_hospital,pathology_comment,ipd_clinical_diagnosis_ppt,tumor_type,sample_type,sample_material,sample_info_comment,tumor_content,ppt_template,output_ppt_file)

        
					print("Generate report for " + DNA_sampleID)
					ppt_nr += 1

	if(ppt_nr > 1):	
		print("Go through the InPreD_PRONTO_metadata file, " + str(ppt_nr) +" reports are generated.")
	else:
		print("Go through the InPreD_PRONTO_metadata file, " + str(ppt_nr) +" report is generated.")

if __name__ == '__main__':
    main(sys.argv[1:])
